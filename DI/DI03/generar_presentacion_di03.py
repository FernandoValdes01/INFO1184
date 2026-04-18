from __future__ import annotations

import os
import shutil
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent
ASSETS_DIR = BASE_DIR / "assets_di03"
OUTPUT_PATH = BASE_DIR / "DI_3_INFO1184_Juan_Munoz.pptx"
MPLCONFIGDIR = ASSETS_DIR / ".mplconfig"
LEGACY_MPLCONFIGDIR = BASE_DIR / ".mplconfig"

ASSETS_DIR.mkdir(exist_ok=True)
MPLCONFIGDIR.mkdir(exist_ok=True)
os.environ.setdefault("MPLCONFIGDIR", str(MPLCONFIGDIR))

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import Circle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace("#", "")
    return RGBColor(
        int(hex_value[0:2], 16),
        int(hex_value[2:4], 16),
        int(hex_value[4:6], 16),
    )


THEME = {
    "paper": rgb("F6F3EE"),
    "white": rgb("FFFFFF"),
    "navy": rgb("11263C"),
    "teal": rgb("2D6F73"),
    "coral": rgb("C7684A"),
    "gold": rgb("B48846"),
    "ink": rgb("1D2730"),
    "muted": rgb("66727F"),
    "line": rgb("D9DDD8"),
    "soft_blue": rgb("EAF1F8"),
    "soft_teal": rgb("EEF6F4"),
    "soft_gold": rgb("FBF3E7"),
    "soft_coral": rgb("F9ECE7"),
    "soft_gray": rgb("F2F4F5"),
    "cluster1": rgb("2967A5"),
    "cluster2": rgb("2F7E67"),
    "noise": rgb("868E96"),
}


COURSE = {
    "code": "INFO1184",
    "name": "Inteligencia de Negocios",
    "semester": "Semestre I-2026",
    "task": "Disertacion 3",
    "topic": "DBSCAN",
    "subtitle": "Density-Based Spatial Clustering of Applications with Noise",
    "team": "Juan Munoz | Vicente Rivera | Fernando Valdes",
}


def ensure_assets_dir() -> None:
    ASSETS_DIR.mkdir(exist_ok=True)


def generate_dataset() -> np.ndarray:
    rng = np.random.default_rng(7)

    t = rng.uniform(0.15, np.pi - 0.05, 70)
    moon = np.column_stack(
        [
            1.3 + 1.75 * np.cos(t) + rng.normal(0, 0.07, t.size),
            0.5 + 1.15 * np.sin(t) + rng.normal(0, 0.07, t.size),
        ]
    )
    blob = rng.normal(loc=[4.2, 3.6], scale=[0.30, 0.26], size=(48, 2))
    noise = rng.uniform(low=[-0.7, -0.5], high=[5.6, 4.8], size=(14, 2))
    return np.vstack([moon, blob, noise])


def dbscan(points: np.ndarray, eps: float, min_pts: int) -> np.ndarray:
    n = len(points)
    labels = np.full(n, -99, dtype=int)
    cluster_id = 0

    def region_query(index: int) -> np.ndarray:
        dist = np.sqrt(((points - points[index]) ** 2).sum(axis=1))
        return np.where(dist <= eps)[0]

    for i in range(n):
        if labels[i] != -99:
            continue
        neighbors = region_query(i)
        if len(neighbors) < min_pts:
            labels[i] = -1
            continue

        labels[i] = cluster_id
        seeds = [int(value) for value in neighbors if value != i]
        cursor = 0
        while cursor < len(seeds):
            j = seeds[cursor]
            if labels[j] == -1:
                labels[j] = cluster_id
            if labels[j] != -99:
                cursor += 1
                continue
            labels[j] = cluster_id
            neighbors_j = region_query(j)
            if len(neighbors_j) >= min_pts:
                for q in neighbors_j:
                    q_int = int(q)
                    if q_int not in seeds:
                        seeds.append(q_int)
            cursor += 1

        cluster_id += 1

    return labels


def classify_points(points: np.ndarray, labels: np.ndarray, eps: float, min_pts: int) -> np.ndarray:
    kinds = np.full(len(points), "border", dtype=object)
    for i in range(len(points)):
        dist = np.sqrt(((points - points[i]) ** 2).sum(axis=1))
        neighbors = np.where(dist <= eps)[0]
        if labels[i] == -1:
            kinds[i] = "noise"
        elif len(neighbors) >= min_pts:
            kinds[i] = "core"
        else:
            kinds[i] = "border"
    return kinds


def generate_figures() -> dict[str, object]:
    ensure_assets_dir()
    points = generate_dataset()
    eps = 0.38
    min_pts = 5
    labels = dbscan(points, eps=eps, min_pts=min_pts)
    kinds = classify_points(points, labels, eps=eps, min_pts=min_pts)

    clusters = sorted(label for label in set(labels.tolist()) if label != -1)
    noise_count = int((labels == -1).sum())

    fig1_path = ASSETS_DIR / "fig01_dbscan_datos.png"
    plt.style.use("default")
    fig, ax = plt.subplots(figsize=(6.2, 4.1), dpi=180)
    ax.scatter(points[:, 0], points[:, 1], s=34, c="#2D6F73", edgecolors="white", linewidths=0.5)
    ax.set_title("Datos sin etiquetar", fontsize=13, color="#11263C", pad=10)
    ax.set_xlabel("X1", fontsize=10)
    ax.set_ylabel("X2", fontsize=10)
    ax.grid(alpha=0.18)
    ax.set_facecolor("#FBFBFA")
    fig.patch.set_facecolor("#FBFBFA")
    fig.tight_layout()
    fig.savefig(fig1_path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)

    fig2_path = ASSETS_DIR / "fig02_dbscan_resultado.png"
    palette = {0: "#2967A5", 1: "#2F7E67", -1: "#868E96"}
    fig, ax = plt.subplots(figsize=(6.2, 4.1), dpi=180)
    for label in sorted(set(labels.tolist())):
        mask = labels == label
        color = palette.get(label, "#C7684A")
        label_name = "Ruido" if label == -1 else f"Cluster {label + 1}"
        ax.scatter(
            points[mask, 0],
            points[mask, 1],
            s=36,
            c=color,
            edgecolors="white",
            linewidths=0.5,
            label=label_name,
        )
    ax.set_title("Resultado con DBSCAN", fontsize=13, color="#11263C", pad=10)
    ax.set_xlabel("X1", fontsize=10)
    ax.set_ylabel("X2", fontsize=10)
    ax.grid(alpha=0.18)
    ax.legend(frameon=False, loc="upper left", fontsize=9)
    ax.set_facecolor("#FBFBFA")
    fig.patch.set_facecolor("#FBFBFA")
    fig.tight_layout()
    fig.savefig(fig2_path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)

    fig3_path = ASSETS_DIR / "fig03_dbscan_concepto.png"
    core_candidates = np.where((labels == 0) & (kinds == "core"))[0]
    core_index = int(core_candidates[8]) if len(core_candidates) > 8 else int(core_candidates[0])
    core_point = points[core_index]
    dists = np.sqrt(((points - core_point) ** 2).sum(axis=1))
    neighborhood = dists <= eps

    fig, ax = plt.subplots(figsize=(5.6, 4.2), dpi=180)
    ax.scatter(points[:, 0], points[:, 1], s=24, c="#D7DEE7", edgecolors="none")
    ax.scatter(points[labels == -1, 0], points[labels == -1, 1], s=28, c="#868E96", label="Noise")
    ax.scatter(
        points[(labels != -1) & (kinds == "border"), 0],
        points[(labels != -1) & (kinds == "border"), 1],
        s=30,
        c="#B48846",
        label="Border",
    )
    ax.scatter(
        points[(labels != -1) & (kinds == "core"), 0],
        points[(labels != -1) & (kinds == "core"), 1],
        s=34,
        c="#2967A5",
        label="Core",
    )
    ax.scatter(core_point[0], core_point[1], s=95, c="#C7684A", edgecolors="white", linewidths=0.8, zorder=5)
    circle = Circle((core_point[0], core_point[1]), eps, fill=False, linestyle="--", linewidth=2, edgecolor="#C7684A")
    ax.add_patch(circle)
    ax.scatter(points[neighborhood, 0], points[neighborhood, 1], s=36, facecolors="none", edgecolors="#C7684A", linewidths=0.8)
    ax.set_title("Vecindad eps y tipos de punto", fontsize=13, color="#11263C", pad=10)
    ax.set_xlabel("X1", fontsize=10)
    ax.set_ylabel("X2", fontsize=10)
    ax.grid(alpha=0.18)
    ax.legend(frameon=False, loc="upper left", fontsize=8)
    ax.set_facecolor("#FBFBFA")
    fig.patch.set_facecolor("#FBFBFA")
    fig.tight_layout()
    fig.savefig(fig3_path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)

    return {
        "points": points,
        "labels": labels,
        "eps": eps,
        "min_pts": min_pts,
        "clusters": len(clusters),
        "noise_count": noise_count,
        "noise_pct": round(noise_count / len(points) * 100, 1),
        "fig_raw": fig1_path,
        "fig_clustered": fig2_path,
        "fig_concept": fig3_path,
    }


def configure_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_shape_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_line(shape, color: RGBColor | None = None, width: Pt = Pt(1.0)) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = width


def add_rect(slide, left, top, width, height, fill, line=None, rounded=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill)
    set_shape_line(shape, line)
    return shape


def style_text_frame(text_frame, margin=0.08, vertical=MSO_ANCHOR.TOP) -> None:
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    text_frame.vertical_anchor = vertical


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    font_name="Aptos",
    size=18,
    color=None,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    fill=None,
    line=None,
    rounded=False,
    margin=0.08,
    vertical=MSO_ANCHOR.TOP,
):
    if fill is None and line is None:
        shape = slide.shapes.add_textbox(left, top, width, height)
    else:
        shape = add_rect(slide, left, top, width, height, fill or THEME["white"], line=line, rounded=rounded)
    style_text_frame(shape.text_frame, margin=margin, vertical=vertical)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or THEME["ink"]
    return shape


def add_paragraphs(
    slide,
    items: list[dict[str, object]],
    left,
    top,
    width,
    height,
    *,
    fill=None,
    line=None,
    rounded=False,
    margin=0.08,
):
    if fill is None and line is None:
        shape = slide.shapes.add_textbox(left, top, width, height)
    else:
        shape = add_rect(slide, left, top, width, height, fill or THEME["white"], line=line, rounded=rounded)
    style_text_frame(shape.text_frame, margin=margin)
    shape.text_frame.clear()
    for index, item in enumerate(items):
        paragraph = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(item.get("space_after", 5))
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = str(item.get("text", ""))
        run.font.name = item.get("font_name", "Aptos")
        run.font.size = Pt(item.get("size", 12))
        run.font.bold = bool(item.get("bold", False))
        run.font.italic = bool(item.get("italic", False))
        run.font.color.rgb = item.get("color", THEME["ink"])
    return shape


def add_badge(slide, text: str, left, top, width, fill, color) -> None:
    add_text(
        slide,
        text,
        left,
        top,
        width,
        Inches(0.30),
        font_name="Aptos",
        size=9.2,
        color=color,
        bold=True,
        fill=fill,
        rounded=True,
        margin=0.04,
    )


def add_footer(slide, number: int) -> None:
    add_text(
        slide,
        f"{COURSE['code']}  |  {COURSE['task']}",
        Inches(0.62),
        Inches(7.02),
        Inches(2.30),
        Inches(0.18),
        font_name="Aptos",
        size=8.5,
        color=THEME["muted"],
        margin=0,
    )
    add_text(
        slide,
        f"{number:02d}",
        Inches(12.32),
        Inches(6.98),
        Inches(0.28),
        Inches(0.22),
        font_name="Aptos",
        size=9,
        color=THEME["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_slide_base(slide, number: int, section: str, title: str, subtitle: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["paper"]
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.16), THEME["navy"])
    add_badge(slide, section, Inches(0.62), Inches(0.34), Inches(1.55), THEME["soft_blue"], THEME["navy"])
    add_text(
        slide,
        title,
        Inches(0.62),
        Inches(0.72),
        Inches(7.60),
        Inches(0.46),
        font_name="Aptos Display",
        size=25,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        subtitle,
        Inches(0.62),
        Inches(1.14),
        Inches(8.10),
        Inches(0.24),
        font_name="Aptos",
        size=11,
        color=THEME["muted"],
        margin=0,
    )
    add_footer(slide, number)


def add_card_title(slide, text: str, left, top, width, color) -> None:
    add_text(
        slide,
        text,
        left,
        top,
        width,
        Inches(0.22),
        font_name="Aptos Display",
        size=15,
        color=color,
        bold=True,
        margin=0,
    )


def add_caption(slide, text: str, left, top, width) -> None:
    add_text(
        slide,
        text,
        left,
        top,
        width,
        Inches(0.24),
        font_name="Aptos",
        size=8.3,
        color=THEME["muted"],
        italic=True,
        margin=0,
    )


def build_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["navy"]

    add_rect(slide, Inches(7.95), Inches(-0.15), Inches(5.6), Inches(7.9), THEME["teal"])
    add_rect(slide, Inches(8.55), Inches(0.45), Inches(4.5), Inches(6.2), THEME["paper"], rounded=True)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.18), THEME["coral"])

    add_badge(slide, f"{COURSE['task']}  |  {COURSE['code']}", Inches(0.72), Inches(0.56), Inches(2.40), THEME["soft_gold"], THEME["gold"])
    add_text(
        slide,
        COURSE["topic"],
        Inches(0.72),
        Inches(1.28),
        Inches(5.10),
        Inches(0.70),
        font_name="Aptos Display",
        size=32,
        color=THEME["white"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Agrupamiento basado en densidad para detectar\nformas arbitrarias y ruido",
        Inches(0.72),
        Inches(2.02),
        Inches(5.80),
        Inches(0.80),
        font_name="Aptos",
        size=17,
        color=THEME["paper"],
        margin=0,
    )
    add_rect(slide, Inches(0.72), Inches(3.05), Inches(4.85), Inches(1.08), THEME["soft_teal"], rounded=True)
    add_text(
        slide,
        "Idea central",
        Inches(0.94),
        Inches(3.26),
        Inches(1.30),
        Inches(0.20),
        font_name="Aptos",
        size=10,
        color=THEME["teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "DBSCAN encuentra clusters sin fijar k y separa automaticamente los puntos atipicos.",
        Inches(0.94),
        Inches(3.52),
        Inches(4.18),
        Inches(0.34),
        font_name="Aptos",
        size=12,
        color=THEME["ink"],
        margin=0,
    )

    add_text(
        slide,
        COURSE["subtitle"],
        Inches(0.72),
        Inches(4.55),
        Inches(5.90),
        Inches(0.28),
        font_name="Aptos",
        size=12.5,
        color=THEME["soft_blue"],
        margin=0,
    )
    add_text(
        slide,
        COURSE["team"],
        Inches(0.72),
        Inches(5.22),
        Inches(5.80),
        Inches(0.28),
        font_name="Aptos",
        size=12,
        color=THEME["white"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        f"{COURSE['name']}  |  {COURSE['semester']}",
        Inches(0.72),
        Inches(5.62),
        Inches(4.70),
        Inches(0.24),
        font_name="Aptos",
        size=10.5,
        color=THEME["soft_blue"],
        margin=0,
    )

    cluster_specs = [
        (9.28, 1.18, 0.23, THEME["cluster1"]),
        (9.75, 1.78, 0.19, THEME["cluster1"]),
        (10.36, 1.52, 0.17, THEME["cluster1"]),
        (10.76, 2.03, 0.21, THEME["cluster1"]),
        (9.98, 2.26, 0.16, THEME["cluster1"]),
        (11.35, 4.10, 0.22, THEME["cluster2"]),
        (11.80, 4.52, 0.18, THEME["cluster2"]),
        (12.18, 4.02, 0.16, THEME["cluster2"]),
        (11.62, 3.56, 0.20, THEME["cluster2"]),
        (12.04, 3.48, 0.18, THEME["cluster2"]),
        (9.46, 5.38, 0.10, THEME["coral"]),
        (10.34, 5.02, 0.08, THEME["coral"]),
        (12.24, 1.76, 0.08, THEME["coral"]),
    ]
    for x, y, radius, color in cluster_specs:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x),
            Inches(y),
            Inches(radius),
            Inches(radius),
        )
        set_shape_fill(shape, color)
        set_shape_line(shape, None)


def build_agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 2, "Panorama", "Temario y propuesta", "Estructura alineada con la pauta general de disertaciones")

    add_rect(slide, Inches(0.62), Inches(1.62), Inches(4.36), Inches(4.98), THEME["white"], line=THEME["line"], rounded=True)
    add_card_title(slide, "Temario", Inches(0.90), Inches(1.86), Inches(1.70), THEME["navy"])
    add_paragraphs(
        slide,
        [
            {"text": "1. Contexto y problema del clustering", "size": 12.2, "color": THEME["ink"], "bold": True},
            {"text": "2. Funcionamiento de DBSCAN: eps, MinPts y expansion", "size": 12.2, "color": THEME["ink"], "bold": True},
            {"text": "3. Ejemplo demostrativo con datos sinteticos", "size": 12.2, "color": THEME["ink"], "bold": True},
            {"text": "4. Aplicaciones del metodo en casos reales", "size": 12.2, "color": THEME["ink"], "bold": True},
            {"text": "5. Ventajas, limitaciones y cierre", "size": 12.2, "color": THEME["ink"], "bold": True},
        ],
        Inches(0.90),
        Inches(2.26),
        Inches(3.70),
        Inches(2.35),
        margin=0,
    )
    add_rect(slide, Inches(0.90), Inches(5.08), Inches(3.52), Inches(0.98), THEME["soft_blue"], rounded=True)
    add_text(
        slide,
        "Proposicion",
        Inches(1.10),
        Inches(5.30),
        Inches(1.10),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "DBSCAN es especialmente util cuando el numero de grupos es desconocido y el ruido no puede ignorarse.",
        Inches(1.10),
        Inches(5.55),
        Inches(3.02),
        Inches(0.26),
        font_name="Aptos",
        size=10.3,
        color=THEME["ink"],
        margin=0,
    )

    add_rect(slide, Inches(5.28), Inches(1.62), Inches(7.42), Inches(1.34), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Gancho",
        Inches(5.62),
        Inches(1.92),
        Inches(0.90),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["soft_teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "En datos reales no siempre existen clusters redondos ni limpios; muchas veces hay formas irregulares y observaciones aisladas.",
        Inches(5.62),
        Inches(2.16),
        Inches(6.56),
        Inches(0.28),
        font_name="Aptos Display",
        size=18,
        color=THEME["white"],
        bold=True,
        margin=0,
    )

    cards = [
        (
            Inches(5.28),
            "Problema",
            "K-means necesita fijar k y suele asumir grupos compactos, por lo que puede forzar particiones poco naturales.",
            THEME["soft_coral"],
            THEME["coral"],
        ),
        (
            Inches(7.82),
            "Respuesta de DBSCAN",
            "Agrupa por densidad local: une puntos densos, reconoce bordes y deja fuera el ruido.",
            THEME["soft_teal"],
            THEME["teal"],
        ),
        (
            Inches(10.36),
            "Valor analitico",
            "Entrega clusters mas cercanos a la forma real de los datos cuando existen outliers o geometrias no esfericas.",
            THEME["soft_gold"],
            THEME["gold"],
        ),
    ]
    for left, title, body, fill, accent in cards:
        add_rect(slide, left, Inches(3.28), Inches(2.18), Inches(2.78), fill, line=THEME["line"], rounded=True)
        add_card_title(slide, title, left + Inches(0.22), Inches(3.55), Inches(1.70), accent)
        add_text(
            slide,
            body,
            left + Inches(0.22),
            Inches(3.98),
            Inches(1.72),
            Inches(1.34),
            font_name="Aptos",
            size=10.4,
            color=THEME["ink"],
            margin=0,
        )


def build_problem_slide(prs: Presentation, visuals: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 3, "Contexto", "Que problema resuelve DBSCAN", "Agrupar cuando la forma de los datos importa tanto como la distancia")

    add_rect(slide, Inches(0.62), Inches(1.62), Inches(4.26), Inches(4.94), THEME["white"], line=THEME["line"], rounded=True)
    add_card_title(slide, "Cuando un metodo tradicional falla", Inches(0.90), Inches(1.88), Inches(2.60), THEME["navy"])
    add_paragraphs(
        slide,
        [
            {"text": "- El numero de clusters puede ser desconocido desde el inicio.", "size": 11.1, "color": THEME["ink"]},
            {"text": "- La forma real puede ser curva, alargada o fragmentada.", "size": 11.1, "color": THEME["ink"]},
            {"text": "- Los puntos atipicos distorsionan centroides y medias.", "size": 11.1, "color": THEME["ink"]},
            {"text": "- En bases reales hay zonas densas y zonas poco pobladas.", "size": 11.1, "color": THEME["ink"]},
        ],
        Inches(0.90),
        Inches(2.34),
        Inches(3.50),
        Inches(1.62),
        margin=0,
    )

    add_rect(slide, Inches(0.90), Inches(4.52), Inches(3.62), Inches(1.40), THEME["soft_blue"], rounded=True)
    add_text(
        slide,
        "Lectura clave",
        Inches(1.12),
        Inches(4.76),
        Inches(1.25),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "DBSCAN no pregunta primero cuantos grupos existen; primero identifica donde realmente hay densidad suficiente.",
        Inches(1.12),
        Inches(5.04),
        Inches(3.02),
        Inches(0.36),
        font_name="Aptos",
        size=10.5,
        color=THEME["ink"],
        margin=0,
    )

    add_rect(slide, Inches(5.14), Inches(1.62), Inches(7.56), Inches(4.94), THEME["white"], line=THEME["line"], rounded=True)
    slide.shapes.add_picture(str(visuals["fig_raw"]), Inches(5.42), Inches(1.98), height=Inches(3.70))
    add_caption(
        slide,
        "Fig. 1. Nube de datos con forma irregular y posibles outliers. Elaboracion propia.",
        Inches(5.58),
        Inches(5.84),
        Inches(6.88),
    )


def build_method_slide(prs: Presentation, visuals: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 4, "Metodo", "Como funciona DBSCAN", "El cluster nace desde vecindades densas y se expande por conectividad")

    add_rect(slide, Inches(0.62), Inches(1.64), Inches(5.05), Inches(4.92), THEME["white"], line=THEME["line"], rounded=True)
    add_card_title(slide, "Conceptos base", Inches(0.90), Inches(1.88), Inches(1.90), THEME["navy"])
    add_badge(slide, "eps", Inches(0.92), Inches(2.28), Inches(0.68), THEME["soft_gold"], THEME["gold"])
    add_text(
        slide,
        "Radio que define la vecindad local de un punto.",
        Inches(1.74),
        Inches(2.26),
        Inches(2.98),
        Inches(0.24),
        font_name="Aptos",
        size=10.8,
        color=THEME["ink"],
        margin=0,
    )
    add_badge(slide, "MinPts", Inches(0.92), Inches(2.72), Inches(1.00), THEME["soft_teal"], THEME["teal"])
    add_text(
        slide,
        "Minimo de puntos dentro de eps para considerar densidad suficiente.",
        Inches(2.00),
        Inches(2.70),
        Inches(2.72),
        Inches(0.36),
        font_name="Aptos",
        size=10.8,
        color=THEME["ink"],
        margin=0,
    )
    add_badge(slide, "Core", Inches(0.92), Inches(3.26), Inches(0.82), THEME["soft_blue"], THEME["cluster1"])
    add_text(
        slide,
        "Punto cuyo vecindario contiene al menos MinPts.",
        Inches(1.86),
        Inches(3.24),
        Inches(2.90),
        Inches(0.24),
        font_name="Aptos",
        size=10.8,
        color=THEME["ink"],
        margin=0,
    )
    add_badge(slide, "Border", Inches(0.92), Inches(3.70), Inches(0.94), THEME["soft_gold"], THEME["gold"])
    add_text(
        slide,
        "Punto cercano a un core, pero sin densidad propia suficiente.",
        Inches(1.98),
        Inches(3.68),
        Inches(2.78),
        Inches(0.36),
        font_name="Aptos",
        size=10.8,
        color=THEME["ink"],
        margin=0,
    )
    add_badge(slide, "Noise", Inches(0.92), Inches(4.22), Inches(0.88), THEME["soft_coral"], THEME["coral"])
    add_text(
        slide,
        "Observacion que no queda conectada a ningun cluster denso.",
        Inches(1.92),
        Inches(4.20),
        Inches(2.82),
        Inches(0.36),
        font_name="Aptos",
        size=10.8,
        color=THEME["ink"],
        margin=0,
    )

    add_text(
        slide,
        "N_eps(p) = { q in D : dist(p, q) <= eps }",
        Inches(0.94),
        Inches(5.04),
        Inches(3.76),
        Inches(0.26),
        font_name="Aptos",
        size=14,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "(1)",
        Inches(4.96),
        Inches(5.04),
        Inches(0.26),
        Inches(0.24),
        font_name="Aptos",
        size=12,
        color=THEME["muted"],
        margin=0,
    )

    add_rect(slide, Inches(5.92), Inches(1.64), Inches(6.78), Inches(2.48), THEME["white"], line=THEME["line"], rounded=True)
    add_card_title(slide, "Secuencia del algoritmo", Inches(6.20), Inches(1.88), Inches(2.30), THEME["navy"])
    steps = [
        ("01", "Seleccionar un punto no visitado."),
        ("02", "Contar su vecindad eps."),
        ("03", "Si cumple MinPts, iniciar o expandir un cluster."),
        ("04", "Si no cumple, marcarlo como ruido provisional."),
    ]
    for idx, (code, text) in enumerate(steps):
        top = 2.30 + idx * 0.40
        add_badge(slide, code, Inches(6.22), Inches(top), Inches(0.52), THEME["navy"], THEME["white"])
        add_text(
            slide,
            text,
            Inches(6.86),
            Inches(top),
            Inches(4.92),
            Inches(0.18),
            font_name="Aptos",
            size=10.8,
            color=THEME["ink"],
            margin=0,
        )

    add_rect(slide, Inches(5.92), Inches(4.38), Inches(6.78), Inches(2.18), THEME["white"], line=THEME["line"], rounded=True)
    slide.shapes.add_picture(str(visuals["fig_concept"]), Inches(6.10), Inches(4.58), height=Inches(1.68))
    add_caption(
        slide,
        "Fig. 2. Punto core, borde y ruido dentro de una vecindad eps. Elaboracion propia.",
        Inches(6.18),
        Inches(6.22),
        Inches(5.90),
    )


def build_example_slide(prs: Presentation, visuals: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 5, "Ejemplo", "Ejemplo demostrativo", "DBSCAN identifica dos zonas densas y deja fuera observaciones aisladas")

    add_rect(slide, Inches(0.62), Inches(1.62), Inches(7.30), Inches(4.94), THEME["white"], line=THEME["line"], rounded=True)
    slide.shapes.add_picture(str(visuals["fig_clustered"]), Inches(0.88), Inches(1.92), height=Inches(3.86))
    add_caption(
        slide,
        "Fig. 3. Resultado de DBSCAN en un conjunto sintetico con forma curva, cluster compacto y ruido. Elaboracion propia.",
        Inches(1.06),
        Inches(5.96),
        Inches(6.46),
    )

    add_rect(slide, Inches(8.16), Inches(1.62), Inches(4.54), Inches(1.34), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Parametros usados",
        Inches(8.46),
        Inches(1.92),
        Inches(1.60),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["soft_teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        f"eps = {visuals['eps']:.2f}  |  MinPts = {visuals['min_pts']}",
        Inches(8.46),
        Inches(2.16),
        Inches(3.18),
        Inches(0.22),
        font_name="Aptos Display",
        size=18,
        color=THEME["white"],
        bold=True,
        margin=0,
    )

    stats = [
        ("Clusters detectados", str(visuals["clusters"]), THEME["soft_blue"], THEME["cluster1"]),
        ("Ruido detectado", f"{visuals['noise_count']} puntos", THEME["soft_teal"], THEME["cluster2"]),
        ("Porcentaje ruido", f"{visuals['noise_pct']}%", THEME["soft_gold"], THEME["gold"]),
    ]
    for idx, (title, value, fill, accent) in enumerate(stats):
        left = 8.16 + (idx % 1) * 0
        top = 3.22 + idx * 0.82
        add_rect(slide, Inches(left), Inches(top), Inches(4.54), Inches(0.62), fill, line=THEME["line"], rounded=True)
        add_text(
            slide,
            title,
            Inches(left + 0.24),
            Inches(top + 0.16),
            Inches(2.30),
            Inches(0.18),
            font_name="Aptos",
            size=10,
            color=accent,
            bold=True,
            margin=0,
        )
        add_text(
            slide,
            value,
            Inches(left + 2.66),
            Inches(top + 0.13),
            Inches(1.42),
            Inches(0.18),
            font_name="Aptos Display",
            size=17,
            color=THEME["navy"],
            bold=True,
            align=PP_ALIGN.RIGHT,
            margin=0,
        )

    add_rect(slide, Inches(8.16), Inches(5.80), Inches(4.54), Inches(0.74), THEME["soft_coral"], rounded=True)
    add_text(
        slide,
        "Interpretacion",
        Inches(8.40),
        Inches(6.00),
        Inches(1.12),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["coral"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "La ventaja principal es que la forma curva se mantiene como un solo cluster sin obligarlo a ser esferico.",
        Inches(9.60),
        Inches(5.98),
        Inches(2.70),
        Inches(0.24),
        font_name="Aptos",
        size=10,
        color=THEME["ink"],
        margin=0,
    )


def build_applications_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 6, "Aplicaciones", "Donde se usa DBSCAN", "El metodo es util cuando la densidad revela patrones mas valiosos que un centroide promedio")

    applications = [
        (
            Inches(0.62),
            "Analitica geoespacial",
            "Detecta zonas densas de accidentes, delitos, demanda de reparto o flujo de clientes.",
            "Funciona bien porque encuentra hotspots con forma irregular y separa eventos aislados.",
            THEME["soft_blue"],
            THEME["cluster1"],
        ),
        (
            Inches(4.52),
            "Fraude y anomalias",
            "Separa patrones normales de transacciones y deja como ruido operaciones atipicas.",
            "Aporta valor porque el outlier no es un error: muchas veces es justamente el caso que interesa investigar.",
            THEME["soft_teal"],
            THEME["cluster2"],
        ),
        (
            Inches(8.42),
            "Sensores y vision",
            "Agrupa puntos en nubes LiDAR, seguimiento de objetos o deteccion de estructuras en imagenes.",
            "Es adecuado porque reconoce objetos conectados por densidad aun cuando sus bordes no sean convexos.",
            THEME["soft_gold"],
            THEME["gold"],
        ),
    ]

    for left, title, desc1, desc2, fill, accent in applications:
        add_rect(slide, left, Inches(1.86), Inches(3.28), Inches(4.54), fill, line=THEME["line"], rounded=True)
        add_text(
            slide,
            title,
            left + Inches(0.28),
            Inches(2.20),
            Inches(2.34),
            Inches(0.24),
            font_name="Aptos Display",
            size=17,
            color=THEME["navy"],
            bold=True,
            margin=0,
        )
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.28), Inches(2.64), Inches(0.28), Inches(0.28))
        set_shape_fill(circle, accent)
        set_shape_line(circle, None)
        add_text(
            slide,
            "Caso de uso",
            left + Inches(0.68),
            Inches(2.66),
            Inches(1.14),
            Inches(0.18),
            font_name="Aptos",
            size=9.6,
            color=accent,
            bold=True,
            margin=0,
        )
        add_text(
            slide,
            desc1,
            left + Inches(0.28),
            Inches(3.06),
            Inches(2.72),
            Inches(0.82),
            font_name="Aptos",
            size=10.5,
            color=THEME["ink"],
            margin=0,
        )
        add_rect(slide, left + Inches(0.28), Inches(4.18), Inches(2.72), Inches(1.10), THEME["white"], line=THEME["line"], rounded=True)
        add_text(
            slide,
            "Por que DBSCAN encaja",
            left + Inches(0.44),
            Inches(4.40),
            Inches(1.62),
            Inches(0.18),
            font_name="Aptos",
            size=9.4,
            color=accent,
            bold=True,
            margin=0,
        )
        add_text(
            slide,
            desc2,
            left + Inches(0.44),
            Inches(4.66),
            Inches(2.38),
            Inches(0.44),
            font_name="Aptos",
            size=9.6,
            color=THEME["ink"],
            margin=0,
        )


def build_strengths_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 7, "Evaluacion", "Ventajas y limitaciones", "Entender cuando usar DBSCAN es tan importante como saber ejecutarlo")

    add_rect(slide, Inches(0.62), Inches(1.72), Inches(5.88), Inches(4.58), THEME["soft_teal"], line=THEME["line"], rounded=True)
    add_card_title(slide, "Ventajas", Inches(0.92), Inches(2.00), Inches(1.30), THEME["teal"])
    add_paragraphs(
        slide,
        [
            {"text": "- No exige definir k antes de comenzar.", "size": 11.5, "color": THEME["ink"]},
            {"text": "- Encuentra clusters con forma arbitraria.", "size": 11.5, "color": THEME["ink"]},
            {"text": "- Identifica ruido y outliers de forma explicita.", "size": 11.5, "color": THEME["ink"]},
            {"text": "- Suele ser interpretable en datos geograficos o espaciales.", "size": 11.5, "color": THEME["ink"]},
        ],
        Inches(0.92),
        Inches(2.46),
        Inches(4.98),
        Inches(2.18),
        margin=0,
    )

    add_rect(slide, Inches(6.82), Inches(1.72), Inches(5.88), Inches(4.58), THEME["soft_coral"], line=THEME["line"], rounded=True)
    add_card_title(slide, "Limitaciones", Inches(7.12), Inches(2.00), Inches(1.60), THEME["coral"])
    add_paragraphs(
        slide,
        [
            {"text": "- Elegir eps y MinPts puede ser sensible al contexto.", "size": 11.5, "color": THEME["ink"]},
            {"text": "- Con densidades muy distintas puede fragmentar o fusionar clusters.", "size": 11.5, "color": THEME["ink"]},
            {"text": "- La escala de las variables afecta directamente la distancia.", "size": 11.5, "color": THEME["ink"]},
            {"text": "- En alta dimension la nocion de densidad pierde fuerza.", "size": 11.5, "color": THEME["ink"]},
        ],
        Inches(7.12),
        Inches(2.46),
        Inches(4.98),
        Inches(2.18),
        margin=0,
    )

    add_rect(slide, Inches(0.62), Inches(6.44), Inches(12.08), Inches(0.52), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Buena practica: escalar variables, probar varios valores de eps y validar el resultado con conocimiento del dominio.",
        Inches(0.92),
        Inches(6.58),
        Inches(11.42),
        Inches(0.18),
        font_name="Aptos",
        size=10.3,
        color=THEME["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )


def build_conclusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 8, "Cierre", "Conclusiones", "DBSCAN es fuerte cuando el ruido y la forma real de los datos no pueden simplificarse")

    add_rect(slide, Inches(0.62), Inches(1.74), Inches(12.08), Inches(1.44), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Conclusion central",
        Inches(0.98),
        Inches(2.02),
        Inches(1.30),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["soft_teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "En dashboards de Inteligencia de Negocios, DBSCAN ayuda a revelar zonas densas, patrones y anomalias que no siempre son evidentes. Esto permite detectar oportunidades, anticipar riesgos y enfocar decisiones sobre los comportamientos que realmente importan.",
        Inches(0.98),
        Inches(2.28),
        Inches(10.78),
        Inches(0.64),
        font_name="Aptos Display",
        size=16,
        color=THEME["white"],
        bold=True,
        margin=0,
    )

    items = [
        (
            Inches(0.62),
            "Cuando conviene usarlo",
            "Cuando existen outliers, clusters no esfericos o el numero de grupos aun no es conocido.",
            THEME["soft_blue"],
            THEME["cluster1"],
        ),
        (
            Inches(4.70),
            "Que debemos cuidar",
            "La seleccion de eps, el escalamiento de variables y la presencia de densidades muy desiguales.",
            THEME["soft_gold"],
            THEME["gold"],
        ),
        (
            Inches(8.78),
            "Mensaje final",
            "No todos los problemas de agrupamiento se resuelven con centroides; a veces la densidad es la verdadera estructura.",
            THEME["soft_teal"],
            THEME["cluster2"],
        ),
    ]
    for left, title, body, fill, accent in items:
        add_rect(slide, left, Inches(3.34), Inches(3.92), Inches(2.76), fill, line=THEME["line"], rounded=True)
        add_text(
            slide,
            title,
            left + Inches(0.24),
            Inches(3.64),
            Inches(2.20),
            Inches(0.22),
            font_name="Aptos Display",
            size=15,
            color=accent,
            bold=True,
            margin=0,
        )
        add_text(
            slide,
            body,
            left + Inches(0.24),
            Inches(4.10),
            Inches(2.26),
            Inches(1.72),
            font_name="Aptos",
            size=10.6,
            color=THEME["ink"],
            margin=0,
        )

    add_text(
        slide,
        "Integrantes: Juan Munoz, Vicente Rivera y Fernando Valdes",
        Inches(0.62),
        Inches(6.68),
        Inches(4.70),
        Inches(0.20),
        font_name="Aptos",
        size=9.5,
        color=THEME["muted"],
        margin=0,
    )


def build_references_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, 9, "Fuentes", "Bibliografia de apoyo", "Referencias base para la definicion, discusion y ejemplo del metodo")

    add_rect(slide, Inches(0.62), Inches(1.70), Inches(12.08), Inches(4.96), THEME["white"], line=THEME["line"], rounded=True)
    refs = [
        "Ester, M., Kriegel, H.-P., Sander, J., & Xu, X. (1996). A density-based algorithm for discovering clusters in large spatial databases with noise. Proceedings of KDD-96, 226-231.",
        "Han, J., Kamber, M., & Pei, J. (2012). Data Mining: Concepts and Techniques (3rd ed.). Morgan Kaufmann.",
        "Schubert, E., Sander, J., Ester, M., Kriegel, H.-P., & Xu, X. (2017). DBSCAN revisited, revisited: Why and how you should still use DBSCAN. ACM Transactions on Database Systems, 42(3), 1-21.",
        "Presentacion elaborada segun la pauta general y rubrica de disertaciones del curso INFO1184, Semestre I-2026.",
    ]
    items = []
    for ref in refs:
        items.append({"text": f"- {ref}", "size": 11.4, "color": THEME["ink"], "space_after": 10})
    add_paragraphs(
        slide,
        items,
        Inches(0.96),
        Inches(2.12),
        Inches(11.10),
        Inches(3.70),
        margin=0,
    )
    add_rect(slide, Inches(0.96), Inches(5.98), Inches(10.72), Inches(0.40), THEME["soft_gray"], rounded=True)
    add_text(
        slide,
        "Las figuras incluidas en la presentacion son de elaboracion propia.",
        Inches(1.22),
        Inches(6.10),
        Inches(10.18),
        Inches(0.18),
        font_name="Aptos",
        size=9.5,
        color=THEME["muted"],
        margin=0,
    )


def build_presentation() -> Path:
    visuals = generate_figures()
    prs = configure_presentation()

    build_title_slide(prs)
    build_agenda_slide(prs)
    build_problem_slide(prs, visuals)
    build_method_slide(prs, visuals)
    build_example_slide(prs, visuals)
    build_applications_slide(prs)
    build_strengths_slide(prs)
    build_conclusion_slide(prs)
    build_references_slide(prs)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


def main() -> None:
    try:
        output_path = build_presentation()
        print(output_path)
    finally:
        shutil.rmtree(MPLCONFIGDIR, ignore_errors=True)
        shutil.rmtree(LEGACY_MPLCONFIGDIR, ignore_errors=True)


if __name__ == "__main__":
    main()
