from __future__ import annotations

import csv
import re
import shutil
import subprocess
from collections import Counter
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parent.parent
TASK_DIR = ROOT_DIR / "TA" / "TA02"
TASK_PDF = TASK_DIR / "Tarea 2 INFO1184.pdf"
DATASET_PATH = TASK_DIR / "data" / "bank-full.csv"
OUTPUT_PATH = BASE_DIR / "TA02_Presentacion_DI02.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace("#", "")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


THEME = {
    "paper": rgb("F7F5F1"),
    "white": rgb("FFFFFF"),
    "navy": rgb("10253E"),
    "slate": rgb("34506B"),
    "teal": rgb("4C7A73"),
    "gold": rgb("B8875B"),
    "ink": rgb("1E2730"),
    "muted": rgb("65727F"),
    "line": rgb("D9DDD7"),
    "sand": rgb("E9E4DC"),
    "soft_blue": rgb("EAF1F8"),
    "soft_teal": rgb("EDF5F2"),
    "soft_gold": rgb("FAF1E7"),
    "soft_risk": rgb("F8ECEA"),
    "risk": rgb("8B554B"),
    "cluster1": rgb("2D5D9F"),
    "cluster2": rgb("4C7A73"),
    "cluster3": rgb("B8875B"),
    "cluster4": rgb("8B554B"),
}


CLUSTER_DATA = [
    {
        "name": "Cluster 1",
        "share": "71,1%",
        "n": "1.422",
        "accent": THEME["cluster1"],
        "fill": THEME["soft_blue"],
        "summary": "Cliente estandar: balance de 1.700 EUR, sin mora ni prestamo personal.",
    },
    {
        "name": "Cluster 2",
        "share": "11,6%",
        "n": "231",
        "accent": THEME["cluster2"],
        "fill": THEME["soft_teal"],
        "summary": "Clientes con prestamo personal activo (100%) y balance reducido (576 EUR).",
    },
    {
        "name": "Cluster 3",
        "share": "15,7%",
        "n": "313",
        "accent": THEME["cluster3"],
        "fill": THEME["soft_gold"],
        "summary": "Segmento con contacto previo: pdays ~ 237 días y previous ~ 3,5.",
    },
    {
        "name": "Cluster 4",
        "share": "1,7%",
        "n": "34",
        "accent": THEME["cluster4"],
        "fill": THEME["soft_risk"],
        "summary": "Grupo de alto riesgo: 100% en mora y balance negativo (-383 EUR).",
    },
]


def clean_text(value: str) -> str:
    if "Ã" in value or "Â" in value:
        try:
            return value.encode("latin1").decode("utf-8")
        except Exception:
            return value
    return value


def extract_course_info(pdf_path: Path) -> dict[str, str]:
    defaults = {
        "course_code": "INFO1184",
        "course_name": "Inteligencia de Negocios",
        "semester": "Semestre I-2026",
        "task": "Tarea 2",
        "subtitle": "Investigación y Desarrollo",
    }
    pdftotext = shutil.which("pdftotext")
    if not pdftotext:
        return defaults

    try:
        result = subprocess.run(
            [pdftotext, "-layout", "-nopgbrk", str(pdf_path), "-"],
            capture_output=True,
            text=True,
            check=True,
        )
    except Exception:
        return defaults

    text = result.stdout
    lines = [clean_text(re.sub(r"\s+", " ", line).strip()) for line in text.splitlines() if line.strip()]

    course_code = defaults["course_code"]
    course_name = defaults["course_name"]
    semester = defaults["semester"]
    task = defaults["task"]
    subtitle = defaults["subtitle"]

    code_match = re.search(r"\b([A-Z]{3,}\d{3,})\b", text)
    if code_match:
        course_code = code_match.group(1)

    for idx, line in enumerate(lines):
        if line == course_code and idx + 1 < len(lines):
            course_name = lines[idx + 1]
            break

    for line in lines:
        if line.lower().startswith("semestre"):
            semester = line
        if line.lower().startswith("tarea"):
            task = line
        if "Investig" in line and "Desarrollo" in line:
            subtitle = line

    return {
        "course_code": clean_text(course_code),
        "course_name": clean_text(course_name),
        "semester": clean_text(semester),
        "task": clean_text(task),
        "subtitle": clean_text(subtitle),
    }


def summarize_bank_dataset(csv_path: Path) -> dict[str, str]:
    job_counter: Counter[str] = Counter()
    total_rows = 0
    yes_rows = 0
    pdays_minus_one = 0
    sum_age = 0.0
    sum_balance = 0.0
    sum_duration = 0.0

    with csv_path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle, delimiter=";")
        fieldnames = reader.fieldnames or []

        for row in reader:
            total_rows += 1
            yes_rows += 1 if row["y"] == "yes" else 0
            pdays_minus_one += 1 if row["pdays"] == "-1" else 0
            sum_age += float(row["age"])
            sum_balance += float(row["balance"])
            sum_duration += float(row["duration"])
            job_counter[row["job"]] += 1

    avg_age = sum_age / total_rows
    avg_balance = sum_balance / total_rows
    avg_duration = sum_duration / total_rows
    yes_pct = yes_rows / total_rows * 100
    no_pct = 100 - yes_pct
    pdays_minus_one_pct = pdays_minus_one / total_rows * 100
    top_jobs = []
    for name, count in job_counter.most_common(3):
        pct = count / total_rows * 100
        top_jobs.append(f"{name} ({pct:.1f}%)")

    return {
        "rows": f"{total_rows:,}".replace(",", "."),
        "columns": str(len(fieldnames)),
        "yes_pct": f"{yes_pct:.1f}".replace(".", ","),
        "no_pct": f"{no_pct:.1f}".replace(".", ","),
        "avg_age": f"{avg_age:.0f}",
        "avg_balance": f"{avg_balance:,.0f}".replace(",", "."),
        "avg_duration": f"{avg_duration:.0f}",
        "pdays_minus_one_pct": f"{pdays_minus_one_pct:.1f}".replace(".", ","),
        "top_jobs": ", ".join(top_jobs),
    }


def configure_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_shape_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_line(shape, color: RGBColor | None = None, width: Pt = Pt(1.0), dash=None) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = width
    if dash is not None:
        shape.line.dash_style = dash


def add_rect(slide, left, top, width, height, fill, line=None, rounded=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill)
    set_shape_line(shape, line)
    return shape


def style_text_frame(text_frame, margin=0.08, vertical=MSO_ANCHOR.TOP) -> None:
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    text_frame.vertical_anchor = vertical


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    font_name="Aptos",
    size=18,
    color=None,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    fill=None,
    line=None,
    rounded=False,
    margin=0.08,
):
    if fill is None and line is None:
        shape = slide.shapes.add_textbox(left, top, width, height)
    else:
        shape = add_rect(slide, left, top, width, height, fill or THEME["white"], line=line, rounded=rounded)
    style_text_frame(shape.text_frame, margin=margin)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or THEME["ink"]
    return shape


def add_paragraphs(
    slide,
    items: list[dict],
    left,
    top,
    width,
    height,
    *,
    fill=None,
    line=None,
    rounded=False,
    margin=0.08,
):
    if fill is None and line is None:
        shape = slide.shapes.add_textbox(left, top, width, height)
    else:
        shape = add_rect(slide, left, top, width, height, fill or THEME["white"], line=line, rounded=rounded)
    style_text_frame(shape.text_frame, margin=margin)
    shape.text_frame.clear()
    for index, item in enumerate(items):
        paragraph = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(item.get("space_after", 4))
        paragraph.space_before = Pt(item.get("space_before", 0))
        run = paragraph.add_run()
        run.text = item["text"]
        run.font.name = item.get("font_name", "Aptos")
        run.font.size = Pt(item.get("size", 18))
        run.font.bold = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.color.rgb = item.get("color", THEME["ink"])
    return shape


def add_picture_fit(slide, image_path: Path, left, top, width, height, *, frame_fill=None, line=None, rounded=False):
    if frame_fill is not None:
        add_rect(slide, left, top, width, height, frame_fill, line=line, rounded=rounded)
        pad = Inches(0.08)
        left += pad
        top += pad
        width -= 2 * pad
        height -= 2 * pad
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        pic_w = width
        pic_h = int(width / img_ratio)
        pic_left = left
        pic_top = int(top + (height - pic_h) / 2)
    else:
        pic_h = height
        pic_w = int(height * img_ratio)
        pic_top = top
        pic_left = int(left + (width - pic_w) / 2)
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_footer(slide, number: int, course_code: str) -> None:
    slide.shapes.add_connector(1, Inches(0.6), Inches(7.05), Inches(12.7), Inches(7.05)).line.color.rgb = THEME["line"]
    add_text(
        slide,
        f"TA02 · {course_code} · Bank Marketing",
        Inches(0.6),
        Inches(7.08),
        Inches(4.5),
        Inches(0.25),
        font_name="Aptos",
        size=8.5,
        color=THEME["muted"],
        margin=0,
    )
    add_text(
        slide,
        f"{number:02d}",
        Inches(12.05),
        Inches(7.02),
        Inches(0.6),
        Inches(0.25),
        font_name="Aptos",
        size=9,
        color=THEME["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_content_scaffold(slide, section: str, title: str, number: int, course_code: str) -> None:
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME["paper"])
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.10), THEME["navy"])
    add_text(
        slide,
        section.upper(),
        Inches(0.6),
        Inches(0.34),
        Inches(2.6),
        Inches(0.22),
        font_name="Aptos",
        size=10,
        color=THEME["teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        title,
        Inches(0.6),
        Inches(0.58),
        Inches(8.6),
        Inches(0.55),
        font_name="Aptos Display",
        size=25,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_footer(slide, number, course_code)


def add_caption(slide, text: str, left, top, width) -> None:
    add_text(
        slide,
        text,
        left,
        top,
        width,
        Inches(0.28),
        font_name="Aptos",
        size=8.5,
        color=THEME["muted"],
        margin=0,
    )


def add_stat_card(slide, value: str, label: str, left, top, width, accent: RGBColor) -> None:
    add_rect(slide, left, top, width, Inches(0.68), THEME["white"], line=THEME["line"], rounded=True)
    add_rect(slide, left, top, Inches(0.08), Inches(0.68), accent)
    add_text(
        slide,
        value,
        left + Inches(0.18),
        top + Inches(0.09),
        width - Inches(0.24),
        Inches(0.25),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        label,
        left + Inches(0.18),
        top + Inches(0.33),
        width - Inches(0.24),
        Inches(0.18),
        font_name="Aptos",
        size=8.5,
        color=THEME["muted"],
        margin=0,
    )


def add_cluster_card(slide, cluster: dict, left, top, width, height) -> None:
    add_rect(slide, left, top, width, height, cluster["fill"], line=THEME["line"], rounded=True)
    add_rect(slide, left, top, Inches(0.10), height, cluster["accent"])
    add_paragraphs(
        slide,
        [
            {"text": cluster["name"], "size": 14, "bold": True, "font_name": "Aptos Display", "color": THEME["navy"], "space_after": 0},
            {"text": f"{cluster['share']} de la muestra · n = {cluster['n']}", "size": 9, "bold": True, "color": cluster["accent"], "space_after": 4},
            {"text": cluster["summary"], "size": 9.5, "color": THEME["ink"], "space_after": 0},
        ],
        left + Inches(0.18),
        top + Inches(0.08),
        width - Inches(0.26),
        height - Inches(0.16),
        margin=0,
    )


def add_action_row(slide, code: str, action: str, left, top, width, accent: RGBColor) -> None:
    add_rect(slide, left, top, width, Inches(0.48), THEME["white"], line=THEME["line"], rounded=True)
    add_rect(slide, left, top, Inches(0.50), Inches(0.48), accent, rounded=True)
    add_text(
        slide,
        code,
        left,
        top + Inches(0.08),
        Inches(0.50),
        Inches(0.18),
        font_name="Aptos",
        size=10,
        color=THEME["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        action,
        left + Inches(0.56),
        top + Inches(0.08),
        width - Inches(0.70),
        Inches(0.20),
        font_name="Aptos",
        size=9.5,
        color=THEME["ink"],
        margin=0,
    )


def build_title_slide(prs: Presentation, course_info: dict[str, str], dataset: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME["paper"])
    add_rect(slide, 0, 0, Inches(4.15), SLIDE_H, THEME["navy"])
    add_rect(slide, Inches(4.15), 0, Inches(0.10), SLIDE_H, THEME["teal"])

    add_text(
        slide,
        f"{course_info['task']} · {course_info['course_code']}",
        Inches(0.6),
        Inches(0.48),
        Inches(2.8),
        Inches(0.22),
        font_name="Aptos",
        size=10,
        color=THEME["sand"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        course_info["course_name"],
        Inches(0.6),
        Inches(0.74),
        Inches(2.8),
        Inches(0.26),
        font_name="Aptos",
        size=12,
        color=THEME["soft_teal"],
        margin=0,
    )
    add_text(
        slide,
        "Análisis de\nAgrupamiento\nJerárquico Aglomerativo",
        Inches(0.6),
        Inches(1.24),
        Inches(3.05),
        Inches(2.10),
        font_name="Aptos Display",
        size=25,
        color=THEME["white"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Dataset Bank Marketing · CRISP-DM · R Project",
        Inches(0.6),
        Inches(3.48),
        Inches(2.9),
        Inches(0.42),
        font_name="Aptos",
        size=12,
        color=THEME["soft_teal"],
        margin=0,
    )
    add_text(
        slide,
        f"{course_info['subtitle']} · {course_info['semester']}",
        Inches(0.6),
        Inches(3.95),
        Inches(2.9),
        Inches(0.26),
        font_name="Aptos",
        size=10.5,
        color=THEME["sand"],
        margin=0,
    )

    add_rect(slide, Inches(4.7), Inches(0.65), Inches(7.95), Inches(4.20), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(
        slide,
        TASK_DIR / "fig_07_clusters_pca.png",
        Inches(4.86),
        Inches(0.82),
        Inches(7.62),
        Inches(3.86),
    )
    add_caption(
        slide,
        "Resultado central del análisis: 4 clusters interpretables proyectados en PCA.",
        Inches(4.86),
        Inches(4.57),
        Inches(7.30),
    )

    add_stat_card(slide, dataset["rows"], "registros del dataset", Inches(4.72), Inches(5.14), Inches(1.85), THEME["navy"])
    add_stat_card(slide, "4", "clusters finales", Inches(6.70), Inches(5.14), Inches(1.55), THEME["teal"])
    add_stat_card(slide, "Power BI", "dashboard requerido", Inches(8.38), Inches(5.14), Inches(2.05), THEME["gold"])
    add_stat_card(slide, dataset["yes_pct"] + "%", "suscribe depósito", Inches(10.56), Inches(5.14), Inches(2.00), THEME["risk"])

    add_rect(slide, Inches(4.72), Inches(6.02), Inches(3.85), Inches(0.94), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Integrantes",
        Inches(4.92),
        Inches(6.16),
        Inches(1.2),
        Inches(0.18),
        font_name="Aptos",
        size=9,
        color=THEME["teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Juan Muñoz\nVicente Rivera\nFernando Valdés",
        Inches(4.92),
        Inches(6.34),
        Inches(2.2),
        Inches(0.44),
        font_name="Aptos",
        size=11,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )

    add_rect(slide, Inches(8.77), Inches(6.02), Inches(3.79), Inches(0.94), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Docente",
        Inches(8.97),
        Inches(6.16),
        Inches(1.0),
        Inches(0.18),
        font_name="Aptos",
        size=9,
        color=THEME["teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Marcos Levano",
        Inches(8.97),
        Inches(6.34),
        Inches(1.8),
        Inches(0.20),
        font_name="Aptos",
        size=11,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Presentacion academica para exposicion final de TA02.",
        Inches(9.95),
        Inches(6.34),
        Inches(2.32),
        Inches(0.28),
        font_name="Aptos",
        size=8.5,
        color=THEME["muted"],
        margin=0,
    )


def build_scope_slide(prs: Presentation, course_info: dict[str, str], dataset: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Introduccion", "Objetivo y alcance de TA02", 2, course_info["course_code"])

    add_paragraphs(
        slide,
        [
            {"text": "La consigna oficial exige aplicar CRISP-DM en sus fases 1, 2, 3 y 4 sobre Bank Marketing, usando agrupamiento aglomerativo jerarquico en R Project y complementando con un dashboard final en Power BI.", "size": 12, "color": THEME["muted"], "space_after": 0},
        ],
        Inches(0.6),
        Inches(1.28),
        Inches(11.8),
        Inches(0.62),
        margin=0,
    )

    add_stat_card(slide, dataset["rows"], "observaciones", Inches(0.6), Inches(1.95), Inches(1.85), THEME["navy"])
    add_stat_card(slide, dataset["columns"], "atributos totales", Inches(2.58), Inches(1.95), Inches(1.90), THEME["teal"])
    add_stat_card(slide, dataset["yes_pct"] + "%", "tasa de suscripcion", Inches(4.61), Inches(1.95), Inches(2.00), THEME["gold"])
    add_stat_card(slide, dataset["pdays_minus_one_pct"] + "%", "sin contacto previo", Inches(6.74), Inches(1.95), Inches(2.20), THEME["risk"])

    cards = [
        ("Caso de negocio", "Segmentar clientes para hacer más eficiente la campaña telefónica de depósitos a plazo.", THEME["navy"], THEME["soft_blue"]),
        ("Método exigido", "Clustering jerárquico aglomerativo con lectura, preparación y modelamiento en R Project.", THEME["teal"], THEME["soft_teal"]),
        ("Entregable visual", "Presentación profesional con resultados, gráficos y una diapositiva reservada para Power BI.", THEME["gold"], THEME["soft_gold"]),
    ]
    card_positions = [Inches(0.6), Inches(4.35), Inches(8.10)]
    for left, (title, body, accent, fill) in zip(card_positions, cards):
        add_rect(slide, left, Inches(2.95), Inches(3.15), Inches(1.45), fill, line=THEME["line"], rounded=True)
        add_rect(slide, left, Inches(2.95), Inches(0.10), Inches(1.45), accent)
        add_text(
            slide,
            title,
            left + Inches(0.16),
            Inches(3.12),
            Inches(2.80),
            Inches(0.24),
            font_name="Aptos Display",
            size=15,
            color=THEME["navy"],
            bold=True,
            margin=0,
        )
        add_text(
            slide,
            body,
            left + Inches(0.16),
            Inches(3.45),
            Inches(2.80),
            Inches(0.70),
            font_name="Aptos",
            size=10.5,
            color=THEME["ink"],
            margin=0,
        )

    roadmap_y = Inches(5.15)
    steps = [
        ("1", "Comprensión\nempresarial", THEME["navy"], THEME["white"]),
        ("2", "Comprensión\nde los datos", THEME["teal"], THEME["white"]),
        ("3", "Preparación\nde datos", THEME["gold"], THEME["white"]),
        ("4", "Modelamiento", THEME["slate"], THEME["white"]),
        ("5", "Dashboard\nPower BI", THEME["sand"], THEME["navy"]),
    ]
    positions = [Inches(0.85), Inches(3.10), Inches(5.35), Inches(7.60), Inches(9.85)]
    for idx, (left, (num, label, fill, text_color)) in enumerate(zip(positions, steps)):
        add_rect(slide, left, roadmap_y, Inches(1.70), Inches(0.88), fill, line=THEME["line"], rounded=True)
        add_text(
            slide,
            f"{num}\n{label}",
            left,
            roadmap_y + Inches(0.09),
            Inches(1.70),
            Inches(0.60),
            font_name="Aptos",
            size=11,
            color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        if idx < len(steps) - 1:
            connector = slide.shapes.add_connector(1, left + Inches(1.72), roadmap_y + Inches(0.44), left + Inches(2.08), roadmap_y + Inches(0.44))
            connector.line.color.rgb = THEME["line"]
            connector.line.width = Pt(2.0)


def build_business_slide(prs: Presentation, course_info: dict[str, str], dataset: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Desarrollo", "Fase 1 · Comprensión empresarial", 3, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.38), Inches(4.00), Inches(4.95), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Problema de negocio",
        Inches(0.82),
        Inches(1.58),
        Inches(2.4),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "- Las campanas telefonicas se realizan sobre una base amplia y heterogenea.", "size": 11, "color": THEME["ink"], "space_after": 2},
            {"text": "- El objetivo es mejorar la focalizacion para aumentar conversion y reducir costo comercial.", "size": 11, "color": THEME["ink"], "space_after": 2},
            {"text": "- Se requiere traducir el comportamiento de clientes en segmentos accionables.", "size": 11, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(0.82),
        Inches(1.95),
        Inches(3.45),
        Inches(1.55),
        margin=0,
    )
    add_text(
        slide,
        "Objetivo de mineria",
        Inches(0.82),
        Inches(3.72),
        Inches(2.2),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "- Identificar grupos homogeneos segun variables demograficas, financieras y de interaccion.", "size": 11, "color": THEME["ink"], "space_after": 2},
            {"text": "- Priorizar clusters interpretables y utiles para decisiones de marketing y riesgo.", "size": 11, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(0.82),
        Inches(4.08),
        Inches(3.45),
        Inches(1.10),
        margin=0,
    )
    add_rect(slide, Inches(0.82), Inches(5.46), Inches(3.45), Inches(0.58), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Criterio de exito: clusters diferenciados, interpretables y utiles para la estrategia comercial.",
        Inches(0.92),
        Inches(5.60),
        Inches(3.20),
        Inches(0.18),
        font_name="Aptos",
        size=9,
        color=THEME["white"],
        bold=True,
        margin=0,
    )

    add_rect(slide, Inches(4.88), Inches(1.38), Inches(3.25), Inches(2.05), THEME["soft_teal"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Preguntas que guian el trabajo",
        Inches(5.10),
        Inches(1.60),
        Inches(2.50),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "1. Existen grupos naturales dentro de la base?", "size": 11, "bold": True, "color": THEME["ink"], "space_after": 3},
            {"text": "2. Que variables definen cada segmento?", "size": 11, "bold": True, "color": THEME["ink"], "space_after": 3},
            {"text": "3. Como usar esos perfiles para orientar campanas?", "size": 11, "bold": True, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(5.10),
        Inches(1.98),
        Inches(2.70),
        Inches(1.12),
        margin=0,
    )

    add_rect(slide, Inches(4.88), Inches(3.63), Inches(3.25), Inches(2.70), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Contexto analítico",
        Inches(5.10),
        Inches(3.86),
        Inches(2.10),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": f"- Dataset: {dataset['rows']} registros y {dataset['columns']} atributos.", "size": 11, "color": THEME["ink"], "space_after": 2},
            {"text": "- Campanas reales de un banco portugues entre 2008 y 2010.", "size": 11, "color": THEME["ink"], "space_after": 2},
            {"text": "- Variable objetivo: suscripcion a deposito a plazo.", "size": 11, "color": THEME["ink"], "space_after": 2},
            {"text": "- Se trabajara una muestra de 2.000 observaciones para el clustering jerarquico.", "size": 11, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(5.10),
        Inches(4.24),
        Inches(2.75),
        Inches(1.55),
        margin=0,
    )

    add_rect(slide, Inches(8.42), Inches(1.38), Inches(4.27), Inches(4.95), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "CRISP-DM en esta tarea",
        Inches(8.66),
        Inches(1.60),
        Inches(2.5),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    phase_y = 2.02
    phases = [
        ("Fase 1", "Comprensión empresarial", THEME["navy"], THEME["white"]),
        ("Fase 2", "Comprensión de los datos", THEME["teal"], THEME["white"]),
        ("Fase 3", "Preparación de los datos", THEME["gold"], THEME["white"]),
        ("Fase 4", "Modelamiento", THEME["slate"], THEME["white"]),
    ]
    for index, (phase, label, fill, txt_color) in enumerate(phases):
        top = Inches(phase_y + index * 0.72)
        add_rect(slide, Inches(8.66), top, Inches(0.95), Inches(0.48), fill, rounded=True)
        add_text(
            slide,
            phase,
            Inches(8.66),
            top + Inches(0.12),
            Inches(0.95),
            Inches(0.15),
            font_name="Aptos",
            size=9,
            color=txt_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(
            slide,
            label,
            Inches(9.78),
            top + Inches(0.12),
            Inches(2.25),
            Inches(0.18),
            font_name="Aptos",
            size=11,
            color=THEME["ink"],
            bold=True,
            margin=0,
        )
    add_rect(slide, Inches(8.66), Inches(5.38), Inches(3.55), Inches(0.58), THEME["soft_gold"], rounded=True)
    add_text(
        slide,
        "La presentación sigue exactamente ese orden: introducción, desarrollo, resultados y cierre.",
        Inches(8.82),
        Inches(5.52),
        Inches(3.20),
        Inches(0.20),
        font_name="Aptos",
        size=9.2,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )


def build_data_slide_one(prs: Presentation, course_info: dict[str, str], dataset: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Desarrollo", "Fase 2 · Comprensión de los datos", 4, course_info["course_code"])

    add_text(
        slide,
        "El analisis exploratorio confirma un fuerte desbalance en la variable objetivo y una alta dispersion financiera, dos elementos clave para la interpretacion posterior de clusters.",
        Inches(0.6),
        Inches(1.26),
        Inches(11.8),
        Inches(0.42),
        font_name="Aptos",
        size=12,
        color=THEME["muted"],
        margin=0,
    )

    add_rect(slide, Inches(0.6), Inches(1.84), Inches(5.70), Inches(3.70), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_01_distribucion_y.png", Inches(0.82), Inches(2.08), Inches(5.26), Inches(3.02))
    add_caption(slide, "Distribución de la variable objetivo y: solo una fracción menor suscribe el depósito.", Inches(0.82), Inches(5.14), Inches(5.10))

    add_rect(slide, Inches(6.48), Inches(1.84), Inches(6.25), Inches(3.70), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_02_edad_balance.png", Inches(6.70), Inches(2.02), Inches(5.82), Inches(3.08))
    add_caption(slide, "Relacion entre edad y balance: dispersion amplia y superposicion entre respuestas yes/no.", Inches(6.70), Inches(5.14), Inches(5.62))

    add_rect(slide, Inches(0.6), Inches(5.78), Inches(12.12), Inches(0.78), THEME["navy"], rounded=True)
    add_text(
        slide,
        f"Insight clave: {dataset['no_pct']}% no suscribe, edad promedio {dataset['avg_age']} anos, balance medio {dataset['avg_balance']} EUR y duracion media del ultimo contacto de {dataset['avg_duration']} segundos.",
        Inches(0.86),
        Inches(5.98),
        Inches(11.60),
        Inches(0.22),
        font_name="Aptos",
        size=10.3,
        color=THEME["white"],
        bold=True,
        margin=0,
    )


def build_data_slide_two(prs: Presentation, course_info: dict[str, str], dataset: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Desarrollo", "Comprensión de los datos · estructura financiera", 5, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.35), Inches(6.25), Inches(4.40), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_03_boxplot_balance_job.png", Inches(0.82), Inches(1.62), Inches(5.82), Inches(3.55))
    add_caption(slide, "Balance por ocupacion: alta dispersion, presencia de outliers y perfiles economicos distintos.", Inches(0.82), Inches(5.26), Inches(5.82))

    add_rect(slide, Inches(7.02), Inches(1.35), Inches(5.70), Inches(4.40), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_04_correlaciones.png", Inches(7.34), Inches(1.68), Inches(5.06), Inches(3.18))
    add_caption(slide, "Matriz de correlaciones entre variables numericas usada para revisar dependencia lineal.", Inches(7.34), Inches(5.02), Inches(4.90))

    add_rect(slide, Inches(0.6), Inches(5.95), Inches(4.05), Inches(0.68), THEME["soft_teal"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        f"Trabajos mas frecuentes: {dataset['top_jobs']}.",
        Inches(0.82),
        Inches(6.12),
        Inches(3.60),
        Inches(0.18),
        font_name="Aptos",
        size=9.6,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )
    add_rect(slide, Inches(4.88), Inches(5.95), Inches(3.60), Inches(0.68), THEME["soft_gold"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        f"{dataset['pdays_minus_one_pct']}% de los clientes no registra contacto previo (pdays = -1).",
        Inches(5.10),
        Inches(6.12),
        Inches(3.16),
        Inches(0.18),
        font_name="Aptos",
        size=9.4,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )
    add_rect(slide, Inches(8.71), Inches(5.95), Inches(4.01), Inches(0.68), THEME["soft_blue"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "No hay una sola variable dominante; la segmentación requiere distancia multivariable y estandarización.",
        Inches(8.93),
        Inches(6.03),
        Inches(3.56),
        Inches(0.32),
        font_name="Aptos",
        size=9.2,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )


def build_preparation_slide(prs: Presentation, course_info: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Desarrollo", "Fase 3 · Preparación de los datos", 6, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.42), Inches(4.00), Inches(4.88), THEME["white"], line=THEME["line"], rounded=True)
    steps = [
        ("1", "Seleccion de variables", "Numericas: age, balance, duration, campaign, pdays y previous."),
        ("2", "Codificacion", "default, housing y loan se binarizan; education se codifica como variable ordinal."),
        ("3", "Muestreo", "Se toma una muestra aleatoria de 2.000 observaciones con semilla 42."),
        ("4", "Escalado", "Se aplica estandarizacion z-score para hacer comparables todas las magnitudes."),
    ]
    top = 1.68
    for number, title, body in steps:
        add_rect(slide, Inches(0.88), Inches(top), Inches(0.52), Inches(0.46), THEME["navy"], rounded=True)
        add_text(
            slide,
            number,
            Inches(0.88),
            Inches(top + 0.10),
            Inches(0.52),
            Inches(0.16),
            font_name="Aptos",
            size=11,
            color=THEME["white"],
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(
            slide,
            title,
            Inches(1.56),
            Inches(top + 0.02),
            Inches(2.70),
            Inches(0.18),
            font_name="Aptos Display",
            size=14,
            color=THEME["navy"],
            bold=True,
            margin=0,
        )
        add_text(
            slide,
            body,
            Inches(1.56),
            Inches(top + 0.22),
            Inches(2.66),
            Inches(0.44),
            font_name="Aptos",
            size=9.6,
            color=THEME["ink"],
            margin=0,
        )
        top += 1.05

    add_rect(slide, Inches(4.92), Inches(1.42), Inches(3.22), Inches(2.05), THEME["soft_blue"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Variables usadas en el modelo",
        Inches(5.16),
        Inches(1.66),
        Inches(2.30),
        Inches(0.22),
        font_name="Aptos Display",
        size=15,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "age · balance · duration · campaign · pdays", "size": 10.5, "bold": True, "color": THEME["ink"], "space_after": 2},
            {"text": "previous · default · housing · loan · education", "size": 10.5, "bold": True, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(5.16),
        Inches(2.08),
        Inches(2.55),
        Inches(0.48),
        margin=0,
    )
    add_text(
        slide,
        "Total: 10 variables de entrada para el clustering.",
        Inches(5.16),
        Inches(2.72),
        Inches(2.40),
        Inches(0.18),
        font_name="Aptos",
        size=9.2,
        color=THEME["muted"],
        margin=0,
    )

    add_rect(slide, Inches(4.92), Inches(3.70), Inches(3.22), Inches(2.60), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Criterio técnico",
        Inches(5.16),
        Inches(3.94),
        Inches(2.00),
        Inches(0.22),
        font_name="Aptos Display",
        size=15,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "- El clustering jerarquico requiere una matriz de distancias n x n.", "size": 10.2, "color": THEME["ink"], "space_after": 2},
            {"text": "- El muestreo reduce costo computacional sin perder interpretabilidad.", "size": 10.2, "color": THEME["ink"], "space_after": 2},
            {"text": "- La estandarizacion evita que balance o duration dominen la distancia euclidiana.", "size": 10.2, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(5.16),
        Inches(4.28),
        Inches(2.62),
        Inches(1.42),
        margin=0,
    )

    add_rect(slide, Inches(8.44), Inches(1.42), Inches(4.28), Inches(4.88), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Formula aplicada",
        Inches(8.82),
        Inches(1.82),
        Inches(1.90),
        Inches(0.22),
        font_name="Aptos",
        size=10,
        color=THEME["sand"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "z = (x - media) / desviación estándar",
        Inches(8.82),
        Inches(2.18),
        Inches(3.20),
        Inches(0.38),
        font_name="Aptos Display",
        size=21,
        color=THEME["white"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "Interpretacion operativa", "size": 12, "bold": True, "color": THEME["soft_teal"], "space_after": 6},
            {"text": "- Cada variable queda en una escala comparable.", "size": 10.5, "color": THEME["white"], "space_after": 3},
            {"text": "- La distancia euclidiana refleja patron global y no solo magnitud absoluta.", "size": 10.5, "color": THEME["white"], "space_after": 3},
            {"text": "- La salida queda lista para el dendrograma y la seleccion de k.", "size": 10.5, "color": THEME["white"], "space_after": 0},
        ],
        Inches(8.82),
        Inches(2.86),
        Inches(3.18),
        Inches(2.15),
        margin=0,
    )


def build_modeling_slide(prs: Presentation, course_info: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Resultados", "Fase 4 · Modelamiento y selección de k", 7, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.38), Inches(6.20), Inches(4.62), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_05_dendrograma_ward.png", Inches(0.82), Inches(1.66), Inches(5.76), Inches(3.78))
    add_caption(slide, "Dendrograma Ward.D2 con corte final en k = 4 clusters.", Inches(0.82), Inches(5.52), Inches(5.76))

    add_rect(slide, Inches(6.96), Inches(1.38), Inches(5.76), Inches(4.62), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_06_elbow_silhouette.png", Inches(7.18), Inches(1.66), Inches(5.32), Inches(3.64))
    add_caption(slide, "Elbow y silueta para k = 2..10: k = 4 logra mejor equilibrio entre forma e interpretación.", Inches(7.18), Inches(5.40), Inches(5.20))

    add_rect(slide, Inches(0.6), Inches(6.14), Inches(12.12), Inches(0.54), THEME["soft_teal"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Se usa distancia euclidiana sobre datos estandarizados y enlace Ward.D2. Aunque k = 3 tiene silueta media mayor (0,2654), se elige k = 4 porque separa un segmento en mora con valor interpretativo y de riesgo.",
        Inches(0.84),
        Inches(6.23),
        Inches(11.62),
        Inches(0.28),
        font_name="Aptos",
        size=9.7,
        color=THEME["ink"],
        bold=True,
        margin=0,
    )


def build_clusters_slide(prs: Presentation, course_info: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Resultados", "Resultado central · cuatro segmentos de clientes", 8, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.38), Inches(6.82), Inches(5.05), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_07_clusters_pca.png", Inches(0.84), Inches(1.72), Inches(6.34), Inches(4.16))
    add_caption(slide, "Proyeccion PCA: se aprecia separacion entre grupos y cierto solapamiento en zonas centrales.", Inches(0.84), Inches(5.98), Inches(6.10))

    add_cluster_card(slide, CLUSTER_DATA[0], Inches(7.70), Inches(1.46), Inches(4.72), Inches(1.08))
    add_cluster_card(slide, CLUSTER_DATA[1], Inches(7.70), Inches(2.72), Inches(4.72), Inches(1.08))
    add_cluster_card(slide, CLUSTER_DATA[2], Inches(7.70), Inches(3.98), Inches(4.72), Inches(1.08))
    add_cluster_card(slide, CLUSTER_DATA[3], Inches(7.70), Inches(5.24), Inches(4.72), Inches(1.08))


def build_profiles_slide(prs: Presentation, course_info: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Resultados", "Perfiles comparados y validación final", 9, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.38), Inches(7.10), Inches(4.28), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_08_perfiles_clusters.png", Inches(0.84), Inches(1.70), Inches(6.64), Inches(3.46))
    add_caption(slide, "Comparación de edad, balance, duración y campaña entre clusters.", Inches(0.84), Inches(5.20), Inches(6.40))

    add_rect(slide, Inches(7.90), Inches(1.38), Inches(4.82), Inches(4.28), THEME["white"], line=THEME["line"], rounded=True)
    add_picture_fit(slide, TASK_DIR / "fig_09_silueta_final.png", Inches(8.18), Inches(1.72), Inches(4.26), Inches(3.40))
    add_caption(slide, "Gráfico de silueta de la asignación final con k = 4.", Inches(8.18), Inches(5.20), Inches(4.04))

    add_rect(slide, Inches(0.6), Inches(5.88), Inches(12.12), Inches(0.72), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Lectura ejecutiva: la estructura es moderada pero útil. Cluster 1 concentra la base estándar, Cluster 3 resalta por historial de contacto y Cluster 4, aunque pequeño, aporta la señal de riesgo más clara para decisiones de negocio.",
        Inches(0.88),
        Inches(6.05),
        Inches(11.58),
        Inches(0.24),
        font_name="Aptos",
        size=10,
        color=THEME["white"],
        bold=True,
        margin=0,
    )


def build_conclusions_slide(prs: Presentation, course_info: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Cierre", "Conclusiones e implicancias de negocio", 10, course_info["course_code"])

    add_rect(slide, Inches(0.6), Inches(1.40), Inches(3.65), Inches(4.90), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Hallazgos finales",
        Inches(0.84),
        Inches(1.66),
        Inches(2.40),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "- La base de clientes no es homogénea: emergen 4 perfiles con sentido comercial.", "size": 10.5, "color": THEME["ink"], "space_after": 3},
            {"text": "- La metodología CRISP-DM ordena el proceso desde el problema hasta el modelamiento.", "size": 10.5, "color": THEME["ink"], "space_after": 3},
            {"text": "- El segmento en mora aporta valor por riesgo, aunque reduzca levemente la silueta.", "size": 10.5, "color": THEME["ink"], "space_after": 3},
            {"text": "- El clustering jerárquico resulta útil para segmentación exploratoria y soporte de decisiones.", "size": 10.5, "color": THEME["ink"], "space_after": 0},
        ],
        Inches(0.84),
        Inches(2.05),
        Inches(3.00),
        Inches(2.10),
        margin=0,
    )
    add_rect(slide, Inches(0.84), Inches(4.72), Inches(3.00), Inches(1.18), THEME["soft_gold"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Limitaciones",
        Inches(1.02),
        Inches(4.92),
        Inches(1.40),
        Inches(0.18),
        font_name="Aptos",
        size=9,
        color=THEME["gold"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Muestra de 2.000 casos por costo O(n^2) y silueta media de 0,2302, esperable en datos reales de marketing.",
        Inches(1.02),
        Inches(5.12),
        Inches(2.64),
        Inches(0.48),
        font_name="Aptos",
        size=9.2,
        color=THEME["ink"],
        margin=0,
    )

    add_rect(slide, Inches(4.54), Inches(1.40), Inches(4.12), Inches(4.90), THEME["white"], line=THEME["line"], rounded=True)
    add_text(
        slide,
        "Acciones sugeridas por segmento",
        Inches(4.80),
        Inches(1.66),
        Inches(2.90),
        Inches(0.22),
        font_name="Aptos Display",
        size=16,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )
    add_action_row(slide, "C1", "Campañas masivas de depósito a plazo para clientes estándar.", Inches(4.80), Inches(2.08), Inches(3.60), THEME["cluster1"])
    add_action_row(slide, "C2", "Ofertas de consolidación o cross-sell para clientes con préstamo.", Inches(4.80), Inches(2.74), Inches(3.60), THEME["cluster2"])
    add_action_row(slide, "C3", "Seguimiento personalizado a clientes con historial de contacto.", Inches(4.80), Inches(3.40), Inches(3.60), THEME["cluster3"])
    add_action_row(slide, "C4", "Excluir de captación y derivar a gestión de riesgo o recuperación.", Inches(4.80), Inches(4.06), Inches(3.60), THEME["cluster4"])
    add_rect(slide, Inches(4.80), Inches(4.88), Inches(3.60), Inches(0.70), THEME["navy"], rounded=True)
    add_text(
        slide,
        "La segmentación convierte hallazgos estadísticos en decisiones comerciales concretas.",
        Inches(5.00),
        Inches(5.06),
        Inches(3.18),
        Inches(0.22),
        font_name="Aptos",
        size=9.3,
        color=THEME["white"],
        bold=True,
        margin=0,
    )

    add_rect(slide, Inches(8.96), Inches(1.40), Inches(3.76), Inches(4.90), THEME["navy"], rounded=True)
    add_text(
        slide,
        "Siguiente paso",
        Inches(9.26),
        Inches(1.74),
        Inches(1.60),
        Inches(0.22),
        font_name="Aptos",
        size=10,
        color=THEME["soft_teal"],
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "Dashboard en Power BI",
        Inches(9.26),
        Inches(2.08),
        Inches(2.60),
        Inches(0.36),
        font_name="Aptos Display",
        size=22,
        color=THEME["white"],
        bold=True,
        margin=0,
    )
    add_paragraphs(
        slide,
        [
            {"text": "- Integrar KPIs de suscripción, balance y duration.", "size": 10.2, "color": THEME["white"], "space_after": 3},
            {"text": "- Filtrar por cluster, empleo, educación y estado crediticio.", "size": 10.2, "color": THEME["white"], "space_after": 3},
            {"text": "- Conectar visualmente el análisis exploratorio con la toma de decisiones.", "size": 10.2, "color": THEME["white"], "space_after": 0},
        ],
        Inches(9.26),
        Inches(2.72),
        Inches(2.80),
        Inches(1.20),
        margin=0,
    )
    add_rect(slide, Inches(9.26), Inches(4.56), Inches(2.90), Inches(0.82), THEME["soft_teal"], rounded=True)
    add_text(
        slide,
        "La última diapositiva queda reservada para insertar ese panel sin rehacer el deck.",
        Inches(9.46),
        Inches(4.74),
        Inches(2.48),
        Inches(0.28),
        font_name="Aptos",
        size=9.4,
        color=THEME["navy"],
        bold=True,
        margin=0,
    )


def build_power_bi_placeholder(prs: Presentation, course_info: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_scaffold(slide, "Power BI", "Placeholder para dashboard final", 11, course_info["course_code"])

    add_text(
        slide,
        "Espacio reservado para insertar la captura o exportación del dashboard final requerido por la tarea.",
        Inches(0.6),
        Inches(1.28),
        Inches(11.8),
        Inches(0.38),
        font_name="Aptos",
        size=12,
        color=THEME["muted"],
        margin=0,
    )

    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(1.90),
        Inches(11.40),
        Inches(4.70),
    )
    frame.fill.background()
    frame.line.color.rgb = THEME["slate"]
    frame.line.width = Pt(2.25)
    frame.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    add_text(
        slide,
        "Dashboard final en Power BI",
        Inches(4.05),
        Inches(3.10),
        Inches(5.20),
        Inches(0.34),
        font_name="Aptos Display",
        size=24,
        color=THEME["navy"],
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        "Insertar aquí el panel con KPIs, filtros y visualizaciones finales del caso Bank Marketing.",
        Inches(3.10),
        Inches(3.56),
        Inches(7.10),
        Inches(0.28),
        font_name="Aptos",
        size=12,
        color=THEME["muted"],
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_stat_card(slide, "KPI", "suscripción / balance / duración", Inches(1.35), Inches(5.16), Inches(2.40), THEME["navy"])
    add_stat_card(slide, "Filtros", "cluster / job / education", Inches(5.05), Inches(5.16), Inches(2.30), THEME["teal"])
    add_stat_card(slide, "Vista", "insights accionables", Inches(8.65), Inches(5.16), Inches(2.50), THEME["gold"])


def build_presentation() -> Path:
    course_info = extract_course_info(TASK_PDF)
    dataset = summarize_bank_dataset(DATASET_PATH)
    prs = configure_presentation()

    build_title_slide(prs, course_info, dataset)
    build_scope_slide(prs, course_info, dataset)
    build_business_slide(prs, course_info, dataset)
    build_data_slide_one(prs, course_info, dataset)
    build_data_slide_two(prs, course_info, dataset)
    build_preparation_slide(prs, course_info)
    build_modeling_slide(prs, course_info)
    build_clusters_slide(prs, course_info)
    build_profiles_slide(prs, course_info)
    build_conclusions_slide(prs, course_info)
    build_power_bi_placeholder(prs, course_info)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


def main() -> None:
    output_path = build_presentation()
    print(output_path)


if __name__ == "__main__":
    main()
